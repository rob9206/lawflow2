"""Convert documents between different formats."""

import os
import io
from typing import Literal

ConversionFormat = Literal["pdf", "png", "txt", "md"]


def convert_pptx_to_pdf(input_path: str, output_path: str) -> None:
    """Convert a PowerPoint file to PDF.
    
    Uses ReportLab to generate PDF from extracted slide content.
    """
    from pptx import Presentation
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Image
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_LEFT
    from PIL import Image as PILImage
    
    prs = Presentation(input_path)
    doc = SimpleDocTemplate(output_path, pagesize=letter)
    styles = getSampleStyleSheet()
    
    # Custom styles
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor='#1a1a1a',
        spaceAfter=20,
        alignment=TA_CENTER,
    )
    
    content_style = ParagraphStyle(
        'CustomContent',
        parent=styles['BodyText'],
        fontSize=12,
        textColor='#333333',
        spaceAfter=12,
        alignment=TA_LEFT,
    )
    
    story = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        # Add slide number
        slide_header = Paragraph(f"Slide {slide_num}", styles['Heading2'])
        story.append(slide_header)
        story.append(Spacer(1, 0.2 * inch))
        
        # Extract text from shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    # Check if it looks like a title (first text or large font)
                    is_title = slide_num == 1 or len(text) < 100
                    style = title_style if is_title else content_style
                    
                    # Clean text for PDF
                    text = text.replace('\x00', '').replace('\r', '\n')
                    para = Paragraph(text, style)
                    story.append(para)
                    story.append(Spacer(1, 0.1 * inch))
            
            # Handle tables
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        row_text = " | ".join(cells)
                        para = Paragraph(row_text, content_style)
                        story.append(para)
                story.append(Spacer(1, 0.1 * inch))
        
        # Add page break between slides
        if slide_num < len(prs.slides):
            story.append(PageBreak())
    
    doc.build(story)


def convert_pptx_to_images(input_path: str, output_dir: str) -> list[str]:
    """Convert PowerPoint slides to PNG images.
    
    Returns list of generated image file paths.
    Note: This extracts text and renders it as images (basic implementation).
    For high-fidelity rendering, consider using LibreOffice or similar.
    """
    from pptx import Presentation
    from PIL import Image, ImageDraw, ImageFont
    
    prs = Presentation(input_path)
    os.makedirs(output_dir, exist_ok=True)
    image_paths = []
    
    # Standard slide dimensions (16:9 aspect ratio)
    width, height = 1920, 1080
    
    for slide_num, slide in enumerate(prs.slides, 1):
        # Create blank image
        img = Image.new('RGB', (width, height), color='white')
        draw = ImageDraw.Draw(img)
        
        # Try to use a nice font, fall back to default
        try:
            title_font = ImageFont.truetype("arial.ttf", 60)
            body_font = ImageFont.truetype("arial.ttf", 40)
        except Exception:
            title_font = ImageFont.load_default()
            body_font = ImageFont.load_default()
        
        y_position = 100
        
        # Extract and draw text
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    # Determine if title or body
                    font = title_font if y_position < 200 else body_font
                    
                    # Wrap text to fit
                    lines = text.split('\n')
                    for line in lines:
                        if y_position < height - 100:
                            draw.text((100, y_position), line, fill='black', font=font)
                            y_position += 80 if font == title_font else 60
        
        # Save image
        output_path = os.path.join(output_dir, f"slide_{slide_num:03d}.png")
        img.save(output_path, 'PNG')
        image_paths.append(output_path)
    
    return image_paths


def convert_pptx_to_text(input_path: str, output_path: str) -> None:
    """Convert PowerPoint to plain text file."""
    from pptx import Presentation
    
    prs = Presentation(input_path)
    
    with open(output_path, 'w', encoding='utf-8') as f:
        for slide_num, slide in enumerate(prs.slides, 1):
            f.write(f"=== Slide {slide_num} ===\n\n")
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        f.write(text + "\n\n")
                
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        if any(cells):
                            f.write(" | ".join(cells) + "\n")
                    f.write("\n")
            
            f.write("\n" + "-" * 60 + "\n\n")


def convert_pptx_to_markdown(input_path: str, output_path: str) -> None:
    """Convert PowerPoint to Markdown format."""
    from pptx import Presentation
    
    prs = Presentation(input_path)
    
    with open(output_path, 'w', encoding='utf-8') as f:
        for slide_num, slide in enumerate(prs.slides, 1):
            f.write(f"# Slide {slide_num}\n\n")
            
            title_written = False
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        if not title_written:
                            f.write(f"## {text}\n\n")
                            title_written = True
                        else:
                            f.write(f"{text}\n\n")
                
                if shape.has_table:
                    table = shape.table
                    # Write table header
                    if len(table.rows) > 0:
                        headers = [cell.text.strip() for cell in table.rows[0].cells]
                        f.write("| " + " | ".join(headers) + " |\n")
                        f.write("|" + "|".join(["---"] * len(headers)) + "|\n")
                        
                        # Write table rows
                        for row in table.rows[1:]:
                            cells = [cell.text.strip() for cell in row.cells]
                            f.write("| " + " | ".join(cells) + " |\n")
                        f.write("\n")
            
            f.write("\n---\n\n")


def convert_document(
    input_path: str,
    output_format: ConversionFormat,
    output_path: str | None = None
) -> str:
    """Convert a document to the specified format.
    
    Args:
        input_path: Path to the input file
        output_format: Target format (pdf, png, txt, md)
        output_path: Optional output path. If None, generates based on input_path
    
    Returns:
        Path to the converted file or directory (for images)
    """
    if not os.path.exists(input_path):
        raise FileNotFoundError(f"Input file not found: {input_path}")
    
    input_ext = os.path.splitext(input_path)[1].lower()
    
    if input_ext != ".pptx":
        raise ValueError(f"Unsupported input format: {input_ext}. Currently only .pptx is supported.")
    
    # Generate output path if not provided
    if output_path is None:
        base_name = os.path.splitext(input_path)[0]
        if output_format == "png":
            output_path = f"{base_name}_slides"
        else:
            output_path = f"{base_name}.{output_format}"
    
    # Route to appropriate converter
    if output_format == "pdf":
        convert_pptx_to_pdf(input_path, output_path)
    elif output_format == "png":
        os.makedirs(output_path, exist_ok=True)
        convert_pptx_to_images(input_path, output_path)
    elif output_format == "txt":
        convert_pptx_to_text(input_path, output_path)
    elif output_format == "md":
        convert_pptx_to_markdown(input_path, output_path)
    else:
        raise ValueError(f"Unsupported output format: {output_format}")
    
    return output_path
