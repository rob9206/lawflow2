"""Extract text content from PDF, PPTX, and DOCX files."""

import os
from dataclasses import dataclass


@dataclass
class ExtractedSection:
    """A section of text extracted from a document."""
    content: str
    section_index: int
    heading: str | None = None
    page_or_slide: int | None = None


def extract_pdf(file_path: str) -> list[ExtractedSection]:
    """Extract text from a PDF file, page by page."""
    import pdfplumber

    sections = []
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text and text.strip():
                sections.append(ExtractedSection(
                    content=text.strip(),
                    section_index=i,
                    page_or_slide=i + 1,
                ))
    return sections


def _extract_table_text(table) -> str:
    """Extract text from a PowerPoint table shape."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        if any(cells):
            rows.append(" | ".join(cells))
    return "\n".join(rows)


def extract_pptx(file_path: str) -> list[ExtractedSection]:
    """Extract text from a PowerPoint file, slide by slide."""
    from pptx import Presentation

    sections = []
    prs = Presentation(file_path)
    for i, slide in enumerate(prs.slides):
        parts = []
        title = None
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    if title is None and shape.shape_type is not None:
                        title = text
                    parts.append(text)
            if shape.has_table:
                table_text = _extract_table_text(shape.table)
                if table_text:
                    parts.append(table_text)
        if parts:
            sections.append(ExtractedSection(
                content="\n\n".join(parts),
                section_index=i,
                heading=title,
                page_or_slide=i + 1,
            ))
    return sections


def extract_docx(file_path: str) -> list[ExtractedSection]:
    """Extract text from a Word document, grouped by headings."""
    from docx import Document

    doc = Document(file_path)
    sections = []
    current_parts: list[str] = []
    current_heading: str | None = None
    section_idx = 0

    def flush():
        nonlocal section_idx
        if current_parts:
            sections.append(ExtractedSection(
                content="\n\n".join(current_parts),
                section_index=section_idx,
                heading=current_heading,
            ))
            section_idx += 1

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # Heading styles start new sections
        if para.style and para.style.name and para.style.name.startswith("Heading"):
            flush()
            current_heading = text
            current_parts = [text]
        else:
            current_parts.append(text)

    flush()
    return sections


def extract_document(file_path: str) -> list[ExtractedSection]:
    """Route to the correct extractor based on file extension."""
    ext = os.path.splitext(file_path)[1].lower()
    extractors = {
        ".pdf": extract_pdf,
        ".pptx": extract_pptx,
        ".docx": extract_docx,
    }
    extractor = extractors.get(ext)
    if not extractor:
        raise ValueError(f"Unsupported file type: {ext}")
    return extractor(file_path)


def chunk_sections(sections: list[ExtractedSection], max_tokens: int = 800) -> list[ExtractedSection]:
    """Split overly long sections into smaller chunks.

    Rough estimate: 1 token ~ 4 characters.
    """
    max_chars = max_tokens * 4
    result = []
    global_idx = 0

    for section in sections:
        if len(section.content) <= max_chars:
            result.append(ExtractedSection(
                content=section.content,
                section_index=global_idx,
                heading=section.heading,
                page_or_slide=section.page_or_slide,
            ))
            global_idx += 1
        else:
            # Split by paragraphs, then by sentences if needed
            paragraphs = section.content.split("\n\n")
            current = ""
            for para in paragraphs:
                if len(current) + len(para) + 2 > max_chars and current:
                    result.append(ExtractedSection(
                        content=current.strip(),
                        section_index=global_idx,
                        heading=section.heading,
                        page_or_slide=section.page_or_slide,
                    ))
                    global_idx += 1
                    current = para
                else:
                    current = current + "\n\n" + para if current else para
            if current.strip():
                result.append(ExtractedSection(
                    content=current.strip(),
                    section_index=global_idx,
                    heading=section.heading,
                    page_or_slide=section.page_or_slide,
                ))
                global_idx += 1

    return result
